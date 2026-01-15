#!/usr/bin/env python3
"""
PPT Generator for AirSaas Flash Reports
Uses python-pptx to generate PowerPoint presentations from fetched data.
Uses Systra template when available.
"""

import json
import os
import copy
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE_PATH = os.path.join(BASE_DIR, 'templates', 'ProjectCardAndFollowUp.pptx')

# =============================================================================
# TEMPLATE COLOR PALETTE (from mapping.json _color_palette)
# MUST be respected in all generated content
# =============================================================================
TEMPLATE_COLORS = {
    'section_titles': RGBColor(0x00, 0x3D, 0x4B),  # #003D4B - dark teal
    'content_text': RGBColor(0x00, 0x00, 0x00),    # #000000 - black
    'date_background': RGBColor(0x78, 0x98, 0x9F), # #78989F - gray-teal
    'date_text': RGBColor(0xFF, 0xFF, 0xFF),       # #FFFFFF - white
}

# Colors from AirSaas moods (matching template palette)
MOOD_COLORS = {
    'good': RGBColor(0x03, 0xE2, 0x6B),       # Green - #03E26B
    'issues': RGBColor(0xFF, 0xD4, 0x3B),     # Yellow - #FFD43B
    'complicated': RGBColor(0xFF, 0x92, 0x2B), # Orange - #FF922B
    'blocked': RGBColor(0xFF, 0x0A, 0x55),    # Red - #FF0A55
}

# Status colors (matching template palette)
STATUS_COLORS = {
    'in_progress': RGBColor(0x00, 0x99, 0xFF),  # Blue - #0099FF
    'ongoing_sourcing': RGBColor(0x00, 0x99, 0xFF),
    'finished': RGBColor(0x03, 0xE2, 0x6B),     # Green - #03E26B
    'backlog': RGBColor(0x99, 0x99, 0x99),      # Gray
    'ideation': RGBColor(0xBB, 0x99, 0xFF),     # Purple
    'solution_chosen': RGBColor(0x00, 0xCC, 0xCC), # Teal
}

# Risk colors (matching template palette)
RISK_COLORS = {
    'low': RGBColor(0x03, 0xE2, 0x6B),    # Green - #03E26B
    'medium': RGBColor(0xFF, 0xD4, 0x3B), # Yellow - #FFD43B
    'high': RGBColor(0xFF, 0x0A, 0x55),   # Red - #FF0A55
}

# French labels for moods
MOOD_LABELS = {
    'good': 'Tout va bien',
    'issues': 'Quelques problèmes',
    'complicated': "C'est compliqué",
    'blocked': 'Bloqué',
}

# Track unfilled fields for Data Notes slide
UNFILLED_FIELDS = []

# =============================================================================
# TEMPLATE STYLES CACHE
# Extracted from template elements and applied uniformly to all generated content
# =============================================================================
TEMPLATE_STYLES = {
    'extracted': False,

    # Bullet formatting
    'bullet': {
        'char': '▪',
        'color': None,  # Hex color string e.g., "D32427"
    },

    # Title style (shape 671 - "Project review :")
    'title': {
        'font_name': None,
        'font_size': None,  # In EMUs
        'font_color': None,  # Hex color string
        'bold': None,
        'italic': None,
    },

    # Inline title style (e.g., "Build", "Made :")
    'inline_title': {
        'font_name': None,
        'font_size': None,
        'font_color': None,
        'bold': True,  # Usually bold
        'italic': None,
    },

    # Content/body text style (bullet items)
    'content': {
        'font_name': None,
        'font_size': None,
        'font_color': None,
        'bold': False,
        'italic': None,
    },

    # Date style (shape 681)
    'date': {
        'font_name': None,
        'font_size': None,
        'font_color': None,
        'bold': None,
    },

    # Status/mood text style (shape 678)
    'status': {
        'font_name': None,
        'font_size': None,
        'font_color': None,
        'bold': None,
    },
}

# Legacy alias for backward compatibility
TEMPLATE_BULLET_FORMAT = {
    'extracted': False,
    'bullet_char': '▪',
    'bullet_color': None,
    'font_name': None,
    'font_size': None,
    'font_color': None,
    'font_bold': False,
}

# Template shape name mappings (from analysis)
# Slide 1 shapes and their purpose:
SLIDE1_SHAPES = {
    'title': 'Google Shape;671;p1',           # "Project review :"
    'mood_comment': 'Google Shape;678;p1',    # Top box "xxx" - mood/status comment
    'info': 'Google Shape;677;p1',            # Deployment area, end users, milestones
    'achievements': 'Google Shape;673;p1',    # Bottom left "xxx" - achievements
    'next_steps': 'Google Shape;679;p1',      # Bottom left lower - next steps
    'budget': 'Google Shape;675;p1',          # Right side - Build BAC/Actual/EAC
    'made': 'Google Shape;676;p1',            # Right side top - Made achievements
    'risks': 'Google Shape;680;p1',           # Right side middle - risks/issues
    'date': 'Google Shape;681;p1',            # Date "dd/mm/yy"
}

# Expected shape positions from ProjectCardAndFollowUp.pptx (for verification)
# Format: (x, y) in inches, rounded to 1 decimal
# Updated: 2026-01-14 from template analysis
EXPECTED_SHAPE_POSITIONS = {
    'title': (0.39, 0.17),         # "Project review :" - Google Shape;671;p1
    'date': (8.88, 0.08),          # "dd/mm/yy" - Google Shape;681;p1
    'mood_status': (0.39, 0.98),   # Status/Mood - Google Shape;678;p1
    'scope': (0.35, 1.93),         # SCOPE & MILESTONES - Google Shape;672;p1
    'info': (0.37, 2.15),          # Deployment area (empty) - Google Shape;677;p1
    'achievements': (0.39, 3.16),  # Description - Google Shape;673;p1
    'trends': (0.39, 4.32),        # TRENDS - Google Shape;674;p1
    'next_steps': (0.39, 4.53),    # Next steps - Google Shape;679;p1
    'made': (5.14, 1.0),           # Completed milestones - Google Shape;676;p1
    'risks': (5.14, 2.48),         # Risk level - Google Shape;680;p1
    'budget': (5.13, 4.4),         # Budget BAC/Actual/EAC - Google Shape;675;p1
}

# Tolerance for position matching (in inches)
POSITION_TOLERANCE = 0.3

# =============================================================================
# TEXT FITTING CONFIGURATION
# =============================================================================
# Shape sizes from template (width x height in inches)
# Used to calculate character limits for each field
SHAPE_SIZES = {
    'title': (7.99, 0.43),          # Wide, single line
    'date': (1.07, 0.35),           # Small, single line
    'mood_status': (4.53, 0.96),    # Medium box, 2-3 lines
    'scope': (4.57, 0.98),          # Medium box, 3-4 lines
    'info': (4.53, 0.75),           # Medium box (disabled)
    'achievements': (4.53, 1.12),   # Medium box, 4-5 lines
    'trends': (4.57, 0.99),         # Medium box, 2-3 lines
    'next_steps': (4.53, 0.89),     # Medium box, 3-4 lines
    'made': (4.51, 1.19),           # Medium box, 4-5 lines
    'risks': (4.41, 1.04),          # Medium box, 3-4 lines
    'budget': (4.43, 1.15),         # Medium box, 5-6 lines
}

# Character limits per field (based on shape size and font size)
# Format: (max_chars_per_line, max_lines, total_char_limit)
TEXT_LIMITS = {
    'title': (80, 1, 80),            # Single line title
    'date': (10, 1, 10),             # Date only
    'mood_status': (50, 3, 120),     # 3 lines max
    'scope': (50, 4, 180),           # 4 lines max
    'achievements': (55, 5, 250),    # 5 lines, truncate long descriptions
    'trends': (50, 3, 120),          # 3 lines max
    'next_steps': (50, 4, 180),      # 4 bullet points max
    'made': (50, 5, 200),            # 5 lines max
    'risks': (50, 4, 180),           # 4 lines max
    'budget': (45, 6, 220),          # 6 lines for budget details
}

# Font sizes per field type
FONT_SIZES = {
    'title': 14,
    'date': 8,
    'content': 9,          # Default for most content
    'content_small': 8,    # For dense content
    'content_tiny': 7,     # For very long content that needs to fit
}

# Minimum font size (never go below this)
MIN_FONT_SIZE = 7


def verify_template(template_path):
    """Verify that the template matches expected shape positions.

    Performs comprehensive verification:
    1. Checks all expected shape positions
    2. Identifies new/unexpected shapes
    3. Reports shape drift (position changes)
    4. Validates template structure

    Returns:
        tuple: (is_valid, warnings, detected_positions, full_report)
        - is_valid: True if template is compatible
        - warnings: List of warning messages
        - detected_positions: Dict of detected shape positions
        - full_report: Dict with complete verification data
    """
    warnings = []
    detected_positions = {}
    full_report = {
        'matched': [],
        'drifted': [],
        'missing': [],
        'new_shapes': [],
        'all_shapes': []
    }

    try:
        prs = Presentation(template_path)
    except Exception as e:
        return False, [f"Cannot open template: {e}"], {}, {}

    if len(prs.slides) < 1:
        return False, ["Template has no slides"], {}, {}

    # Analyze slide 0 (Project Review template)
    slide = prs.slides[0]

    # Build comprehensive map of ALL text shapes with their positions
    shape_positions = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            x = round(shape.left.inches, 2)
            y = round(shape.top.inches, 2)
            w = round(shape.width.inches, 2)
            h = round(shape.height.inches, 2)
            text_preview = shape.text[:50].replace('\n', '↵') if shape.text else "(empty)"
            shape_info = {
                'name': shape.name,
                'x': x,
                'y': y,
                'w': w,
                'h': h,
                'text': text_preview,
                'matched_to': None
            }
            shape_positions.append(shape_info)
            full_report['all_shapes'].append(shape_info.copy())

    # Check each expected position
    matched_count = 0
    for field_name, (expected_x, expected_y) in EXPECTED_SHAPE_POSITIONS.items():
        # Find closest shape to expected position
        best_match = None
        best_distance = float('inf')

        for sp in shape_positions:
            if sp['matched_to'] is not None:
                continue  # Already matched to another field
            dist = ((sp['x'] - expected_x) ** 2 + (sp['y'] - expected_y) ** 2) ** 0.5
            if dist < best_distance:
                best_distance = dist
                best_match = sp

        if best_match and best_distance <= POSITION_TOLERANCE:
            detected_positions[field_name] = (best_match['x'], best_match['y'])
            best_match['matched_to'] = field_name
            matched_count += 1

            match_info = {
                'field': field_name,
                'expected': (expected_x, expected_y),
                'actual': (best_match['x'], best_match['y']),
                'offset': best_distance,
                'shape_name': best_match['name'],
                'text': best_match['text']
            }

            # Categorize match quality
            if best_distance < 0.05:
                full_report['matched'].append(match_info)
            else:
                full_report['drifted'].append(match_info)
                warnings.append(
                    f"  ⚠️  {field_name}: DRIFTED from ({expected_x}, {expected_y}) "
                    f"to ({best_match['x']}, {best_match['y']}) - offset: {best_distance:.2f}\""
                )
        else:
            detected_positions[field_name] = None
            full_report['missing'].append({
                'field': field_name,
                'expected': (expected_x, expected_y)
            })
            warnings.append(
                f"  ❌ {field_name}: Expected at ({expected_x}, {expected_y}) - NOT FOUND"
            )

    # Identify NEW/UNEXPECTED shapes (not in expected positions)
    for sp in shape_positions:
        if sp['matched_to'] is None:
            full_report['new_shapes'].append({
                'name': sp['name'],
                'position': (sp['x'], sp['y']),
                'size': (sp['w'], sp['h']),
                'text': sp['text']
            })
            # Only warn if shape has meaningful content
            if sp['text'] and sp['text'] != '(empty)' and 'xxx' not in sp['text'].lower():
                warnings.append(
                    f"  ℹ️  NEW SHAPE: {sp['name']} at ({sp['x']}, {sp['y']}) - \"{sp['text'][:30]}\""
                )

    # Determine if valid (at least 70% of shapes found)
    match_ratio = matched_count / len(EXPECTED_SHAPE_POSITIONS)
    is_valid = match_ratio >= 0.7

    # Build summary
    summary_lines = [
        f"\n{'─'*60}",
        f"TEMPLATE VERIFICATION REPORT",
        f"{'─'*60}",
        f"Template: {os.path.basename(template_path)}",
        f"Expected shapes: {len(EXPECTED_SHAPE_POSITIONS)}",
        f"Matched perfectly: {len(full_report['matched'])}",
        f"Position drifted: {len(full_report['drifted'])}",
        f"Missing: {len(full_report['missing'])}",
        f"New/unexpected: {len(full_report['new_shapes'])}",
        f"Match ratio: {match_ratio:.0%}",
        f"{'─'*60}"
    ]
    warnings.insert(0, '\n'.join(summary_lines))

    if not is_valid:
        warnings.append(
            f"\n⚠️  Template may have changed significantly! Only {match_ratio:.0%} compatible."
        )
        warnings.append(
            "   Run: python3 scripts/generate_ppt.py --analyze"
        )
        warnings.append(
            "   Then update EXPECTED_SHAPE_POSITIONS in the script."
        )

    if full_report['drifted']:
        warnings.append("\n📍 DRIFTED SHAPES - Update these positions in EXPECTED_SHAPE_POSITIONS:")
        for d in full_report['drifted']:
            warnings.append(f"    '{d['field']}': ({d['actual'][0]}, {d['actual'][1]}),  # was ({d['expected'][0]}, {d['expected'][1]})")

    if full_report['new_shapes']:
        warnings.append(f"\n🆕 NEW SHAPES DETECTED ({len(full_report['new_shapes'])}):")
        for ns in full_report['new_shapes'][:5]:  # Show first 5
            warnings.append(f"    {ns['name']} at ({ns['position'][0]}, {ns['position'][1]}) - \"{ns['text'][:30]}\"")
        if len(full_report['new_shapes']) > 5:
            warnings.append(f"    ... and {len(full_report['new_shapes']) - 5} more")

    return is_valid, warnings, detected_positions, full_report


def analyze_template(template_path, verbose=False):
    """Analyze template and print ALL shape details comprehensively.

    Use this to update EXPECTED_SHAPE_POSITIONS after template changes.

    Args:
        template_path: Path to the template file
        verbose: If True, show even more detail (fill colors, borders, etc.)
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(template_path)

    print(f"\n{'='*80}")
    print(f"COMPREHENSIVE TEMPLATE ANALYSIS")
    print(f"{'='*80}")
    print(f"File: {os.path.basename(template_path)}")
    print(f"Full path: {template_path}")
    print(f"Dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\" (inches)")
    print(f"           {prs.slide_width.emu} x {prs.slide_height.emu} (EMUs)")
    print(f"Total Slides: {len(prs.slides)}")
    print(f"Total Layouts: {len(prs.slide_layouts)}")

    # Analyze layouts
    print(f"\n{'─'*80}")
    print("AVAILABLE LAYOUTS:")
    print(f"{'─'*80}")
    for idx, layout in enumerate(prs.slide_layouts):
        print(f"  [{idx}] {layout.name}")

    # Track all shapes for summary
    all_shapes_summary = []
    text_shapes_by_slide = {}
    image_count = 0
    shape_type_counts = {}

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n{'='*80}")
        print(f"SLIDE {slide_idx}")
        print(f"{'='*80}")
        print(f"Layout: {slide.slide_layout.name}")
        print(f"Total shapes: {len(slide.shapes)}")

        text_shapes_by_slide[slide_idx] = []

        # Categorize shapes
        text_shapes = []
        image_shapes = []
        other_shapes = []

        for shape in slide.shapes:
            shape_type = type(shape).__name__
            shape_type_counts[shape_type] = shape_type_counts.get(shape_type, 0) + 1

            x = round(shape.left.inches, 2)
            y = round(shape.top.inches, 2)
            w = round(shape.width.inches, 2)
            h = round(shape.height.inches, 2)

            shape_info = {
                'name': shape.name,
                'type': shape_type,
                'x': x, 'y': y,
                'w': w, 'h': h,
                'has_text': shape.has_text_frame,
                'is_placeholder': shape.is_placeholder if hasattr(shape, 'is_placeholder') else False,
                'text': '',
                'slide': slide_idx
            }

            if shape.has_text_frame:
                text = shape.text[:100].replace('\n', '↵') if shape.text else "(empty)"
                shape_info['text'] = text
                text_shapes.append(shape_info)
                text_shapes_by_slide[slide_idx].append(shape_info)
            elif hasattr(shape, 'image'):
                image_shapes.append(shape_info)
                image_count += 1
            else:
                other_shapes.append(shape_info)

            all_shapes_summary.append(shape_info)

        # Print text shapes (most important for mapping)
        if text_shapes:
            print(f"\n  📝 TEXT SHAPES ({len(text_shapes)}):")
            print(f"  {'─'*76}")
            for s in sorted(text_shapes, key=lambda x: (x['y'], x['x'])):
                placeholder_mark = " [PLACEHOLDER]" if s['is_placeholder'] else ""
                print(f"  ▸ {s['name']}{placeholder_mark}")
                print(f"    Position: ({s['x']}, {s['y']}) inches")
                print(f"    Size: {s['w']}\" x {s['h']}\"")
                print(f"    Text: \"{s['text'][:60]}{'...' if len(s['text']) > 60 else ''}\"")
                print()

        # Print image shapes
        if image_shapes:
            print(f"\n  🖼️  IMAGE SHAPES ({len(image_shapes)}):")
            print(f"  {'─'*76}")
            for s in image_shapes:
                print(f"  ▸ {s['name']}")
                print(f"    Position: ({s['x']}, {s['y']}) | Size: {s['w']}\" x {s['h']}\"")

        # Print other shapes (if verbose)
        if other_shapes and verbose:
            print(f"\n  🔷 OTHER SHAPES ({len(other_shapes)}):")
            print(f"  {'─'*76}")
            for s in other_shapes:
                print(f"  ▸ {s['name']} ({s['type']})")
                print(f"    Position: ({s['x']}, {s['y']}) | Size: {s['w']}\" x {s['h']}\"")

    # Summary section
    print(f"\n{'='*80}")
    print("SUMMARY")
    print(f"{'='*80}")

    print(f"\nShape type distribution:")
    for stype, count in sorted(shape_type_counts.items(), key=lambda x: -x[1]):
        print(f"  {stype}: {count}")

    print(f"\nTotal images: {image_count}")
    print(f"Total text shapes: {sum(len(shapes) for shapes in text_shapes_by_slide.values())}")

    # Generate EXPECTED_SHAPE_POSITIONS format
    print(f"\n{'='*80}")
    print("SUGGESTED EXPECTED_SHAPE_POSITIONS (for generate_ppt.py)")
    print(f"{'='*80}")
    print("# Copy this to update EXPECTED_SHAPE_POSITIONS constant:")
    print("EXPECTED_SHAPE_POSITIONS = {")

    if 0 in text_shapes_by_slide:
        for s in sorted(text_shapes_by_slide[0], key=lambda x: (x['y'], x['x'])):
            # Generate a field name from the shape name or position
            field_name = s['name'].replace('Google Shape;', 'shape_').replace(';p1', '')
            text_preview = s['text'][:30].replace('↵', ' ') if s['text'] else 'empty'
            print(f"    '{field_name}': ({s['x']}, {s['y']}),  # \"{text_preview}\"")

    print("}")

    # Field mapping suggestions
    print(f"\n{'='*80}")
    print("FIELD MAPPING SUGGESTIONS (Slide 0 - Project Template)")
    print(f"{'='*80}")
    if 0 in text_shapes_by_slide:
        print("\nDetected fields based on position and content:\n")
        for s in sorted(text_shapes_by_slide[0], key=lambda x: (x['y'], x['x'])):
            # Try to identify field purpose
            text_lower = s['text'].lower()
            field_guess = "unknown"

            if 'project review' in text_lower or 'title' in s['name'].lower():
                field_guess = "TITLE"
            elif 'dd/mm' in text_lower or 'date' in text_lower:
                field_guess = "DATE"
            elif 'status' in text_lower or 'mood' in text_lower:
                field_guess = "MOOD_STATUS"
            elif 'scope' in text_lower or 'milestone' in text_lower:
                field_guess = "SCOPE_MILESTONES"
            elif 'description' in text_lower or 'benefit' in text_lower:
                field_guess = "ACHIEVEMENTS"
            elif 'trend' in text_lower or 'progress' in text_lower:
                field_guess = "TRENDS"
            elif 'next' in text_lower or 'step' in text_lower:
                field_guess = "NEXT_STEPS"
            elif 'made' in text_lower or 'done' in text_lower or 'completed' in text_lower:
                field_guess = "MADE"
            elif 'risk' in text_lower or 'issue' in text_lower:
                field_guess = "RISKS"
            elif 'budget' in text_lower or 'bac' in text_lower or 'eac' in text_lower:
                field_guess = "BUDGET"
            elif 'xxx' in text_lower or s['text'] == '(empty)':
                field_guess = "PLACEHOLDER (needs content)"

            print(f"  Position ({s['x']}, {s['y']}): {field_guess}")
            print(f"    Shape: {s['name']}")
            print(f"    Current text: \"{s['text'][:50]}\"")
            print(f"    Size: {s['w']}\" x {s['h']}\"")
            print()

    print(f"{'='*80}")
    print("END OF ANALYSIS")
    print(f"{'='*80}\n")


def export_template_mapping(template_path):
    """Export ALL template shape positions to JSON for documentation.

    Creates a comprehensive JSON file with all shape details that can be
    used for field mapping and documentation.
    """
    prs = Presentation(template_path)

    export_data = {
        '_description': 'Template shape positions extracted by generate_ppt.py',
        '_template': os.path.basename(template_path),
        '_exported_at': datetime.now().isoformat(),
        '_dimensions': {
            'width_inches': round(prs.slide_width.inches, 2),
            'height_inches': round(prs.slide_height.inches, 2)
        },
        'layouts': [],
        'slides': []
    }

    # Export layouts
    for idx, layout in enumerate(prs.slide_layouts):
        export_data['layouts'].append({
            'index': idx,
            'name': layout.name
        })

    # Export slides
    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            'index': slide_idx,
            'layout': slide.slide_layout.name,
            'shapes': []
        }

        for shape in slide.shapes:
            shape_data = {
                'name': shape.name,
                'type': type(shape).__name__,
                'position': {
                    'x': round(shape.left.inches, 2),
                    'y': round(shape.top.inches, 2)
                },
                'size': {
                    'width': round(shape.width.inches, 2),
                    'height': round(shape.height.inches, 2)
                },
                'has_text_frame': shape.has_text_frame,
                'is_placeholder': shape.is_placeholder if hasattr(shape, 'is_placeholder') else False
            }

            if shape.has_text_frame:
                shape_data['text_preview'] = shape.text[:100] if shape.text else ""
                shape_data['paragraph_count'] = len(shape.text_frame.paragraphs)

            slide_data['shapes'].append(shape_data)

        # Sort shapes by position (top to bottom, left to right)
        slide_data['shapes'].sort(key=lambda s: (s['position']['y'], s['position']['x']))
        export_data['slides'].append(slide_data)

    # Save to config directory
    output_path = os.path.join(BASE_DIR, 'config', 'template_shapes.json')
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(export_data, f, indent=2, ensure_ascii=False)

    print(f"✓ Template mapping exported to: {output_path}")
    print(f"  Slides: {len(export_data['slides'])}")
    print(f"  Total shapes: {sum(len(s['shapes']) for s in export_data['slides'])}")

    # Also print summary
    print("\nSummary by slide:")
    for slide in export_data['slides']:
        text_shapes = [s for s in slide['shapes'] if s['has_text_frame']]
        print(f"  Slide {slide['index']} ({slide['layout']}): {len(text_shapes)} text shapes")

    return output_path


# =============================================================================
# TEXT FITTING FUNCTIONS
# =============================================================================

def truncate_text(text, max_chars, ellipsis='...'):
    """Truncate text to max_chars, adding ellipsis if truncated.

    Args:
        text: The text to truncate
        max_chars: Maximum number of characters
        ellipsis: String to append if truncated (default '...')

    Returns:
        Truncated text with ellipsis if needed
    """
    if not text or len(text) <= max_chars:
        return text or ''

    # Find a good break point (word boundary)
    truncated = text[:max_chars - len(ellipsis)]

    # Try to break at last space
    last_space = truncated.rfind(' ')
    if last_space > max_chars * 0.7:  # Only if we don't lose too much
        truncated = truncated[:last_space]

    return truncated.rstrip() + ellipsis


def truncate_lines(text, max_lines, max_chars_per_line=None):
    """Truncate text to a maximum number of lines.

    Args:
        text: The text to truncate (may contain newlines)
        max_lines: Maximum number of lines to keep
        max_chars_per_line: Optional limit per line

    Returns:
        Text truncated to max_lines
    """
    if not text:
        return ''

    lines = text.split('\n')

    # Truncate each line if needed
    if max_chars_per_line:
        lines = [truncate_text(line, max_chars_per_line, '...') if len(line) > max_chars_per_line else line
                 for line in lines]

    # Limit number of lines
    if len(lines) > max_lines:
        lines = lines[:max_lines]
        # Only add ellipsis if last line doesn't already have one
        if lines and lines[-1] and not lines[-1].endswith('...'):
            lines[-1] = lines[-1].rstrip()
            if len(lines[-1]) > 3:
                lines[-1] = lines[-1][:-3] + '...'
            else:
                lines.append('...')

    return '\n'.join(lines)


def fit_text_to_shape(text, field_name, add_prefix_newline=False):
    """Fit text to a shape based on predefined limits.

    Args:
        text: The text to fit
        field_name: Name of the field (must be in TEXT_LIMITS)
        add_prefix_newline: Add newline at start (for shapes with visual titles)

    Returns:
        tuple: (fitted_text, recommended_font_size)
    """
    if not text:
        return ('\n' if add_prefix_newline else '', FONT_SIZES.get('content', 9))

    # Get limits for this field
    limits = TEXT_LIMITS.get(field_name, (50, 4, 180))
    max_chars_per_line, max_lines, total_limit = limits

    # First, apply total character limit
    if len(text) > total_limit:
        text = truncate_text(text, total_limit)

    # Then apply line limits
    text = truncate_lines(text, max_lines, max_chars_per_line)

    # Determine font size based on content length
    font_size = FONT_SIZES.get('content', 9)

    # Reduce font for longer content
    if len(text) > total_limit * 0.8:
        font_size = FONT_SIZES.get('content_small', 8)
    if len(text) > total_limit * 0.95:
        font_size = FONT_SIZES.get('content_tiny', 7)

    # Add prefix newline if needed (for shapes with visual section titles)
    if add_prefix_newline:
        text = '\n' + text

    return (text, font_size)


def format_bullet_list(items, max_items=3, max_chars_per_item=45, prefix='• '):
    """Format a list of items as bullet points with truncation.

    Args:
        items: List of items (strings or dicts with 'name'/'title' key)
        max_items: Maximum number of items to show
        max_chars_per_item: Maximum characters per bullet point
        prefix: Prefix for each bullet (default '• ')

    Returns:
        Formatted bullet list as string
    """
    if not items:
        return ''

    lines = []
    for item in items[:max_items]:
        # Handle both strings and dicts
        if isinstance(item, dict):
            text = item.get('name') or item.get('title') or str(item)
        else:
            text = str(item)

        # Truncate item text
        text = truncate_text(text, max_chars_per_item - len(prefix))
        lines.append(f"{prefix}{text}")

    # Add indicator if there are more items
    if len(items) > max_items:
        remaining = len(items) - max_items
        lines.append(f"  (+{remaining} more)")

    return '\n'.join(lines)


def format_milestone_name(milestone, max_chars=40):
    """Format a milestone name, optionally including date.

    Args:
        milestone: Dict with 'name' and optionally 'date' keys
        max_chars: Maximum total characters

    Returns:
        Formatted milestone string
    """
    name = milestone.get('name', '')
    date = milestone.get('date', '')

    if date:
        # Reserve space for date
        date_str = f" ({format_date(date)})"
        name_limit = max_chars - len(date_str)
        name = truncate_text(name, name_limit)
        return f"{name}{date_str}"
    else:
        return truncate_text(name, max_chars)


def load_data(data_path):
    """Load fetched project data from JSON file."""
    with open(data_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def find_shape_by_name(slide, name):
    """Find a shape by its name."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_shape_text(shape, text, font_size=None, style_type=None):
    """Set text on a shape, completely replacing existing content.

    Applies template styles automatically if style_type is specified.

    Args:
        shape: The shape to modify
        text: The text to set
        font_size: Optional font size in points (e.g., 9 for Pt(9))
        style_type: Optional template style to apply ('title', 'content', 'date', 'status')
    """
    if shape and shape.has_text_frame:
        tf = shape.text_frame
        text_str = str(text) if text else ""

        # Clear ALL existing text by clearing each paragraph
        # We need to keep the first paragraph (can't delete it) but clear others
        while len(tf.paragraphs) > 1:
            # Remove paragraphs by clearing the XML
            p = tf.paragraphs[-1]._p
            tf._txBody.remove(p)

        # Clear the first (and now only) paragraph
        if tf.paragraphs:
            tf.paragraphs[0].clear()

        # Now set the new text - just assign directly
        # This handles newlines as soft returns within the shape
        if tf.paragraphs:
            tf.paragraphs[0].text = text_str

            # Apply template style if specified
            if style_type:
                apply_template_style(tf.paragraphs[0], style_type, override_size=font_size)
            elif font_size is not None:
                # Fallback to just setting font size
                tf.paragraphs[0].font.size = Pt(font_size)


def set_shape_text_with_structure(shape, items, title_text=None, font_size=None, use_bullets=True, style_type='content'):
    """Set text on a shape preserving paragraph structure and bullets.

    This function properly populates shapes that have multiple paragraphs
    with bullet points, preserving the template's visual structure.
    Applies template styles automatically.

    Args:
        shape: The shape to modify
        items: List of strings (one per bullet/paragraph) or single string
        title_text: Optional title for first paragraph (no bullet)
        font_size: Optional font size in points (overrides template)
        use_bullets: Whether to format items as bullet list
        style_type: Template style to apply ('content', 'status', etc.)
    """
    if not shape or not shape.has_text_frame:
        return

    tf = shape.text_frame

    # Convert single string to list of items
    if isinstance(items, str):
        items = [line.strip() for line in items.split('\n') if line.strip()]

    # If no items, just set empty or title only
    if not items and not title_text:
        set_shape_text(shape, '', font_size)
        return

    # Clear existing paragraphs except first
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        tf._txBody.remove(p)

    # Clear first paragraph
    if tf.paragraphs:
        tf.paragraphs[0].clear()

    # Set up content index
    content_idx = 0
    para_idx = 0

    # Handle title in first paragraph if specified
    if title_text:
        p = tf.paragraphs[0]
        p.text = title_text
        # Apply inline_title style from template
        apply_template_style(p, 'inline_title', override_size=font_size)
        # Remove bullet from title paragraph
        _remove_bullet(p)
        para_idx = 1
    else:
        # First item goes in first paragraph
        if items:
            p = tf.paragraphs[0]
            p.text = items[0]
            # Apply content style from template
            apply_template_style(p, style_type, override_size=font_size)
            if use_bullets:
                _add_bullet(p)
            content_idx = 1
            para_idx = 1

    # Add remaining items as new paragraphs
    for item in items[content_idx:]:
        p = tf.add_paragraph()
        p.text = item
        # Apply content style from template
        apply_template_style(p, style_type, override_size=font_size)
        if use_bullets:
            _add_bullet(p)


def _extract_font_color_from_run(run):
    """Extract font color from a run element."""
    try:
        from pptx.oxml.ns import qn
        rPr = run._r.find(qn('a:rPr'))
        if rPr is not None:
            solidFill = rPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    return srgbClr.get('val')
    except:
        pass
    return None


def _extract_paragraph_style(para, style_dict):
    """Extract font style from a paragraph into a style dictionary."""
    if para.runs:
        run = para.runs[0]
        if run.font.size:
            style_dict['font_size'] = run.font.size
        if run.font.name:
            style_dict['font_name'] = run.font.name
        if run.font.bold is not None:
            style_dict['bold'] = run.font.bold
        if run.font.italic is not None:
            style_dict['italic'] = run.font.italic

        # Extract font color
        color = _extract_font_color_from_run(run)
        if color:
            style_dict['font_color'] = color


def load_style_extraction_rules():
    """Load style extraction rules from mapping.json.

    Returns the rules dict or defaults if file not found.
    """
    mapping_path = os.path.join(BASE_DIR, 'config', 'mapping.json')
    try:
        with open(mapping_path, 'r', encoding='utf-8') as f:
            mapping = json.load(f)
            rules = mapping.get('_style_extraction_rules', {})
            if rules:
                return rules
    except Exception as e:
        print(f"  [Warning] Could not load style rules from mapping.json: {e}")

    # Return default rules
    return {
        'style_types': {
            'title': {'extract_from_position': {'x': 0.39, 'y': 0.17}},
            'date': {'extract_from_position': {'x': 8.88, 'y': 0.08}},
            'status': {'extract_from_position': {'x': 0.39, 'y': 0.98}},
            'inline_title': {'extract_from_positions': [
                {'x': 5.14, 'y': 1.0},
                {'x': 5.13, 'y': 4.4}
            ]},
            'content': {'extract_from_y_positions': [1.93, 3.16, 4.53, 2.48]},
            'bullet': {'skip_colors': ['000000', 'FFFFFF']}
        },
        'defaults': {
            'font_name': 'Lato',
            'font_size_pt': 9,
            'font_color': '000000',
            'bullet_char': '▪',
            'bullet_color': 'D32427'
        }
    }


def extract_all_template_styles(slide):
    """Extract all visual styles from the template slide.

    Reads extraction rules from mapping.json._style_extraction_rules
    to determine which shapes to extract styles from.

    This captures styles for:
    - Title, Date, Status (by position from mapping)
    - Inline titles ("Build", "Made :")
    - Content/bullet items
    - Bullet formatting (character, color)

    These styles are then applied uniformly to all generated content.
    """
    global TEMPLATE_STYLES, TEMPLATE_BULLET_FORMAT

    if TEMPLATE_STYLES['extracted']:
        return

    from pptx.oxml.ns import qn

    # Load rules from mapping.json
    rules = load_style_extraction_rules()
    style_types = rules.get('style_types', {})
    defaults = rules.get('defaults', {})

    print("  Scanning template shapes for styles (rules from mapping.json)...")

    # Build shape_purposes map from rules
    shape_purposes = {}

    # Single position extractions
    for style_name in ['title', 'date', 'status']:
        if style_name in style_types:
            pos = style_types[style_name].get('extract_from_position', {})
            if pos:
                shape_purposes[(pos.get('x', 0), pos.get('y', 0))] = style_name

    # Multiple position extractions (inline_title)
    if 'inline_title' in style_types:
        positions = style_types['inline_title'].get('extract_from_positions', [])
        for pos in positions:
            shape_purposes[(pos.get('x', 0), pos.get('y', 0))] = 'inline_title'

    # Content shapes (by Y position)
    bullet_shapes_y = style_types.get('content', {}).get('extract_from_y_positions', [1.93, 3.16, 4.53, 2.48])

    # Bullet skip colors
    bullet_skip_colors = style_types.get('bullet', {}).get('skip_colors', ['000000', 'FFFFFF'])

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # Get shape position
        x = shape.left / 914400  # EMUs to inches
        y = shape.top / 914400

        tf = shape.text_frame
        if not tf.paragraphs:
            continue

        # Check if this shape matches a known purpose by position
        shape_purpose = None
        for (px, py), purpose in shape_purposes.items():
            if abs(x - px) < 0.1 and abs(y - py) < 0.1:
                shape_purpose = purpose
                break

        # Extract styles based on shape purpose
        for para in tf.paragraphs:
            # Check for bullet formatting first
            try:
                pPr = para._p.find(qn('a:pPr'))
                if pPr is not None:
                    # Look for bullet character (only capture first one)
                    if TEMPLATE_STYLES['bullet']['char'] == '▪':  # Default value
                        buChar = pPr.find(qn('a:buChar'))
                        if buChar is not None:
                            char = buChar.get('char')
                            if char:
                                TEMPLATE_STYLES['bullet']['char'] = char
                                TEMPLATE_BULLET_FORMAT['bullet_char'] = char

                    # Look for bullet color (only capture first non-skipped color)
                    # Skip colors are defined in mapping.json._style_extraction_rules.style_types.bullet.skip_colors
                    buClr = pPr.find(qn('a:buClr'))
                    if buClr is not None:
                        srgbClr = buClr.find(qn('a:srgbClr'))
                        if srgbClr is not None:
                            val = srgbClr.get('val')
                            # Skip colors from mapping rules (default: black, white)
                            skip_colors_upper = [c.upper() for c in bullet_skip_colors]
                            if val and val.upper() not in skip_colors_upper:
                                if TEMPLATE_STYLES['bullet']['color'] is None:
                                    TEMPLATE_STYLES['bullet']['color'] = val
                                    TEMPLATE_BULLET_FORMAT['bullet_color'] = val
                                    print(f"    Bullet color: #{val} (captured)")
            except:
                pass

            # Extract font style based on shape purpose
            if shape_purpose == 'title' and TEMPLATE_STYLES['title']['font_size'] is None:
                _extract_paragraph_style(para, TEMPLATE_STYLES['title'])
                if TEMPLATE_STYLES['title']['font_size']:
                    size_pt = TEMPLATE_STYLES['title']['font_size'].pt if hasattr(TEMPLATE_STYLES['title']['font_size'], 'pt') else TEMPLATE_STYLES['title']['font_size'] / 12700
                    print(f"    Title: {TEMPLATE_STYLES['title']['font_name']}, {size_pt:.0f}pt, color=#{TEMPLATE_STYLES['title']['font_color']}")

            elif shape_purpose == 'date' and TEMPLATE_STYLES['date']['font_size'] is None:
                _extract_paragraph_style(para, TEMPLATE_STYLES['date'])
                if TEMPLATE_STYLES['date']['font_size']:
                    print(f"    Date style extracted")

            elif shape_purpose == 'status' and TEMPLATE_STYLES['status']['font_size'] is None:
                _extract_paragraph_style(para, TEMPLATE_STYLES['status'])
                if TEMPLATE_STYLES['status']['font_size']:
                    print(f"    Status style extracted")

            elif shape_purpose == 'inline_title':
                # First paragraph is usually the inline title
                if para == tf.paragraphs[0] and TEMPLATE_STYLES['inline_title']['font_size'] is None:
                    _extract_paragraph_style(para, TEMPLATE_STYLES['inline_title'])
                    if TEMPLATE_STYLES['inline_title']['font_size']:
                        print(f"    Inline title style extracted")
                # Second paragraph would be content
                elif para != tf.paragraphs[0] and TEMPLATE_STYLES['content']['font_size'] is None:
                    _extract_paragraph_style(para, TEMPLATE_STYLES['content'])

            # Check if this is a bullet content shape
            elif any(abs(y - by) < 0.1 for by in bullet_shapes_y):
                if TEMPLATE_STYLES['content']['font_size'] is None:
                    _extract_paragraph_style(para, TEMPLATE_STYLES['content'])
                    if TEMPLATE_STYLES['content']['font_size']:
                        size_pt = TEMPLATE_STYLES['content']['font_size'].pt if hasattr(TEMPLATE_STYLES['content']['font_size'], 'pt') else TEMPLATE_STYLES['content']['font_size'] / 12700
                        print(f"    Content: {TEMPLATE_STYLES['content']['font_name']}, {size_pt:.0f}pt, color=#{TEMPLATE_STYLES['content']['font_color']}")

    # Apply defaults from mapping.json for any styles not extracted
    if TEMPLATE_STYLES['bullet']['color'] is None:
        TEMPLATE_STYLES['bullet']['color'] = defaults.get('bullet_color', 'D32427')
        TEMPLATE_BULLET_FORMAT['bullet_color'] = TEMPLATE_STYLES['bullet']['color']
        print(f"    [Default] Bullet color: #{TEMPLATE_STYLES['bullet']['color']}")

    if TEMPLATE_STYLES['bullet']['char'] == '▪':  # Still default
        default_char = defaults.get('bullet_char', '▪')
        TEMPLATE_STYLES['bullet']['char'] = default_char
        TEMPLATE_BULLET_FORMAT['bullet_char'] = default_char

    if TEMPLATE_STYLES['content']['font_name'] is None:
        TEMPLATE_STYLES['content']['font_name'] = defaults.get('font_name', 'Lato')
        print(f"    [Default] Font: {TEMPLATE_STYLES['content']['font_name']}")

    if TEMPLATE_STYLES['content']['font_color'] is None:
        TEMPLATE_STYLES['content']['font_color'] = defaults.get('font_color', '000000')

    # Mark as extracted and sync legacy format
    TEMPLATE_STYLES['extracted'] = True
    TEMPLATE_BULLET_FORMAT['extracted'] = True
    TEMPLATE_BULLET_FORMAT['font_name'] = TEMPLATE_STYLES['content'].get('font_name')
    TEMPLATE_BULLET_FORMAT['font_size'] = TEMPLATE_STYLES['content'].get('font_size')
    TEMPLATE_BULLET_FORMAT['font_color'] = TEMPLATE_STYLES['content'].get('font_color')
    TEMPLATE_BULLET_FORMAT['font_bold'] = TEMPLATE_STYLES['content'].get('bold', False)

    print(f"  ✓ Template styles extracted (rules from mapping.json)")


def apply_template_style(paragraph, style_type='content', override_size=None):
    """Apply extracted template style to a paragraph.

    Args:
        paragraph: The paragraph to style
        style_type: One of 'title', 'inline_title', 'content', 'date', 'status'
        override_size: Optional font size in points to override template size
    """
    global TEMPLATE_STYLES

    if not TEMPLATE_STYLES['extracted']:
        return

    style = TEMPLATE_STYLES.get(style_type, TEMPLATE_STYLES['content'])

    # Apply font name
    if style.get('font_name'):
        paragraph.font.name = style['font_name']

    # Apply font size (override takes precedence)
    if override_size:
        paragraph.font.size = Pt(override_size)
    elif style.get('font_size'):
        paragraph.font.size = style['font_size']

    # Apply font color
    if style.get('font_color'):
        try:
            color_hex = style['font_color']
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            paragraph.font.color.rgb = RGBColor(r, g, b)
        except:
            pass

    # Apply bold/italic
    if style.get('bold') is not None:
        paragraph.font.bold = style['bold']
    if style.get('italic') is not None:
        paragraph.font.italic = style['italic']


def extract_bullet_format_from_shape(shape):
    """Legacy function - redirects to extract_all_template_styles.

    Kept for backward compatibility. The new extraction is done via
    extract_all_template_styles() which captures all element styles.
    """
    # This function is now a no-op since we use the comprehensive extractor
    pass


def _add_bullet(paragraph, apply_template_style=True):
    """Add bullet point formatting to a paragraph.

    If template style was extracted, applies the same bullet color and formatting.
    """
    global TEMPLATE_BULLET_FORMAT

    try:
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml
        from lxml import etree

        pPr = paragraph._p.get_or_add_pPr()

        # Remove any existing bullet settings
        for child in list(pPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag.startswith('bu'):
                pPr.remove(child)

        # Get bullet character from template or default
        bullet_char = TEMPLATE_BULLET_FORMAT.get('bullet_char', '▪')

        # Add bullet character
        buChar = parse_xml(
            f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="{bullet_char}"/>'
        )
        pPr.append(buChar)

        # Apply bullet color if extracted from template
        if apply_template_style and TEMPLATE_BULLET_FORMAT.get('bullet_color'):
            color_hex = TEMPLATE_BULLET_FORMAT['bullet_color']
            buClr = parse_xml(
                f'<a:buClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:srgbClr val="{color_hex}"/>'
                f'</a:buClr>'
            )
            pPr.append(buClr)

        # Apply font color to match template bullet style
        if apply_template_style and TEMPLATE_BULLET_FORMAT.get('font_color'):
            color_hex = TEMPLATE_BULLET_FORMAT['font_color']
            try:
                r = int(color_hex[0:2], 16)
                g = int(color_hex[2:4], 16)
                b = int(color_hex[4:6], 16)
                paragraph.font.color.rgb = RGBColor(r, g, b)
            except:
                pass

    except Exception as e:
        # Fallback: prepend bullet character to text
        bullet_char = TEMPLATE_BULLET_FORMAT.get('bullet_char', '▪')
        if not paragraph.text.startswith(bullet_char) and not paragraph.text.startswith('•'):
            paragraph.text = f'{bullet_char} ' + paragraph.text


def _remove_bullet(paragraph):
    """Remove bullet point formatting from a paragraph."""
    try:
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml

        pPr = paragraph._p.get_or_add_pPr()

        # Remove any existing bullet settings
        for child in list(pPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag.startswith('bu'):
                pPr.remove(child)

        # Add buNone to explicitly disable bullets
        buNone = parse_xml(
            '<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
        pPr.append(buNone)
    except Exception:
        pass  # Silent fail - title will just not have bullet removed


def set_shape_text_with_title_padding(shape, items, padding_lines=1, font_size=None, use_bullets=True, style_type='content'):
    """Set text on a shape with empty padding lines at top.

    Use this when the layout has a title that overlaps the top of the shape.
    The padding_lines creates empty space so content appears below the title.
    Applies template styles automatically.

    Args:
        shape: The shape to modify
        items: List of strings (one per bullet/paragraph)
        padding_lines: Number of empty lines at top (default 1)
        font_size: Optional font size in points (overrides template)
        use_bullets: Whether to format items as bullet list
        style_type: Template style to apply ('content', 'status', etc.)
    """
    if not shape or not shape.has_text_frame:
        return

    tf = shape.text_frame

    # Convert single string to list of items
    if isinstance(items, str):
        items = [line.strip() for line in items.split('\n') if line.strip()]

    # Clear existing paragraphs except first
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        tf._txBody.remove(p)

    # Clear first paragraph
    if tf.paragraphs:
        tf.paragraphs[0].clear()

    # Add padding lines (empty paragraphs) first
    # First paragraph is already there, set it empty
    p = tf.paragraphs[0]
    p.text = ""
    _remove_bullet(p)  # No bullet for empty line
    if font_size:
        p.font.size = Pt(font_size)

    # Add additional padding lines if needed
    for _ in range(padding_lines - 1):
        p = tf.add_paragraph()
        p.text = ""
        _remove_bullet(p)
        if font_size:
            p.font.size = Pt(font_size)

    # Now add the actual content items with template styles
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        # Apply content style from template
        apply_template_style(p, style_type, override_size=font_size)
        if use_bullets:
            _add_bullet(p)
        else:
            _remove_bullet(p)


def format_date(date_str):
    """Format date string to dd/mm/yyyy."""
    if not date_str:
        return datetime.now().strftime("%d/%m/%Y")
    try:
        dt = datetime.fromisoformat(date_str.replace('Z', '+00:00'))
        return dt.strftime("%d/%m/%Y")
    except:
        return date_str


def duplicate_slide(prs, slide_index):
    """Duplicate a slide including all its content (shapes, images, formatting).

    Uses internal pptx XML manipulation to properly clone slides.
    """
    from copy import deepcopy
    from pptx.parts.slide import SlidePart
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    source_slide = prs.slides[slide_index]

    # Get the slide part and make a copy
    slide_layout = source_slide.slide_layout
    sld_id_lst = prs.slides._sldIdLst

    # Add a new slide using the same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Clear the new slide's shapes (they come from the layout)
    # We'll copy shapes from the source instead
    spTree = new_slide.shapes._spTree

    # Remove all existing shapes except the first element (usually nvGrpSpPr)
    shapes_to_remove = []
    for child in spTree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ['sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp']:
            shapes_to_remove.append(child)

    for shape_el in shapes_to_remove:
        spTree.remove(shape_el)

    # Copy shapes from source slide
    source_spTree = source_slide.shapes._spTree
    for child in source_spTree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ['sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp']:
            new_el = deepcopy(child)
            spTree.append(new_el)

    # Copy image relationships
    source_part = source_slide.part
    new_part = new_slide.part

    for rel in source_part.rels.values():
        if rel.reltype == RT.IMAGE:
            # Add the same image relationship to new slide
            new_part.relate_to(rel._target, RT.IMAGE)

    return new_slide


def populate_project_slide(slide, proj_data, reference_data):
    """Populate a project review slide with data.

    This function finds shapes by position (more reliable than name after duplication)
    and updates their text content.
    """
    project = proj_data.get('project', {})
    resolved = proj_data.get('resolved', {})
    milestones = proj_data.get('milestones', [])
    decisions = proj_data.get('decisions', [])
    attention_points = proj_data.get('attention_points', [])

    # Build a map of shapes by approximate position (in inches, rounded)
    # This is more reliable than names which may change after duplication
    shape_map = {}
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Create a position key (rounded to 1 decimal)
            pos_key = (round(shape.left.inches, 1), round(shape.top.inches, 1))
            shape_map[pos_key] = shape

    # Position-based shape mapping (from template analysis):
    # Title: (0.4, 0.2) - "Project review :"
    # Date: (8.9, 0.1) - "dd/mm/yy"
    # Status/Mood: (0.4, 0.9) - Status info
    # Info: (0.4, 2.1) - Deployment/dates
    # Achievements: (0.4, 3.0) - Description
    # Next steps: (0.4, 4.4) - Next steps
    # Made: (5.1, 0.9) - Completed items
    # Risks: (5.1, 2.3) - Risks
    # Budget: (5.1, 4.2) - Budget info

    # Helper to find shape by approximate position
    # Uses smaller tolerance and finds closest match
    def find_shape_by_pos(target_x, target_y, tolerance=0.15):
        best_match = None
        best_distance = float('inf')

        for (x, y), shape in shape_map.items():
            dist = ((x - target_x) ** 2 + (y - target_y) ** 2) ** 0.5
            if dist <= tolerance and dist < best_distance:
                best_match = shape
                best_distance = dist

        return best_match

    # =========================================================================
    # TITLE - "Project review : {project_name}"
    # Uses 'title' style from template
    # =========================================================================
    title_shape = find_shape_by_pos(0.39, 0.17)
    if title_shape:
        project_name = project.get('name', 'Unknown Project')
        # Truncate long project names to fit in title area
        title_text = f"Project review : {truncate_text(project_name, 60)}"
        set_shape_text(title_shape, title_text, font_size=FONT_SIZES['title'], style_type='title')

    # =========================================================================
    # DATE - User preference: dd/mm/yyyy format
    # Uses 'date' style from template
    # =========================================================================
    date_shape = find_shape_by_pos(8.88, 0.08)
    if date_shape:
        set_shape_text(date_shape, datetime.now().strftime("%d/%m/%Y"), font_size=FONT_SIZES['date'], style_type='date')

    # =========================================================================
    # MOOD/STATUS - User preference: Status + Mood only
    # Uses 'status' style from template
    # =========================================================================
    mood_shape = find_shape_by_pos(0.39, 0.98)
    if mood_shape:
        mood = resolved.get('mood', project.get('mood', ''))
        status = resolved.get('status', project.get('status', ''))

        # Truncate long status values
        status_text = truncate_text(status, 40) if status else 'N/A'
        mood_text = truncate_text(mood, 40) if mood else 'N/A'

        comment = f"Status: {status_text}\nMood: {mood_text}"
        fitted_text, font_size = fit_text_to_shape(comment, 'mood_status')
        set_shape_text(mood_shape, fitted_text, font_size=font_size, style_type='status')

    # =========================================================================
    # SCOPE & MILESTONES - Shape 672
    # CRITICAL: Layout title "SCOPE & MILESTONES" overlaps top of this shape!
    #   - Layout title Y: 1.833" - 2.089"
    #   - Shape 672 Y: 1.931" - 2.906"
    #   - OVERLAP: 0.158" (title covers top 0.16" of shape)
    # SOLUTION: First paragraph must be EMPTY to push content below title
    # NOTE: Shape 677 (info area) OVERLAPS with 672 - DO NOT USE 677
    # =========================================================================
    scope_shape = find_shape_by_pos(0.35, 1.93)
    if scope_shape:
        ms_count = len(milestones) if milestones else 0
        completed_ms = len([m for m in milestones if m.get('status') == 'done']) if milestones else 0
        progress = project.get('progress', 0)
        start_date = project.get('start_date', '')
        end_date = project.get('end_date', '')

        # Build items list - includes dates since we're NOT using shape 677
        scope_items = [
            f"Milestones: {completed_ms}/{ms_count}",
            f"Progress: {progress}%",
        ]
        if start_date:
            scope_items.append(f"Start: {format_date(start_date)}")
        if end_date:
            scope_items.append(f"End: {format_date(end_date)}")

        # Use special function with padding for overlapping title
        set_shape_text_with_title_padding(
            scope_shape,
            scope_items,
            padding_lines=1,  # Empty first line to avoid title overlap
            font_size=FONT_SIZES['content'],
            use_bullets=True
        )

    # =========================================================================
    # INFO AREA (Shape 677) - CLEAR IT (overlaps with 672)
    # This shape overlaps with shape 672 (SCOPE & MILESTONES)
    # All info is now consolidated into shape 672 above
    # We must CLEAR this shape to avoid showing template placeholder text
    # =========================================================================
    info_shape = find_shape_by_pos(0.37, 2.15)
    if info_shape:
        set_shape_text(info_shape, '', font_size=FONT_SIZES['content'])

    # =========================================================================
    # ACHIEVEMENTS - User preference: Description del proyecto
    # NOTE: Layout title "DESCRIPTION & EXPECTED BENEFITS" is at (0.47, 0.63)
    #       Content shape at (0.39, 3.16) - but this area is labeled "PROGRESS" in layout
    #       Template uses bullet formatting (▪)
    # =========================================================================
    achievements_shape = find_shape_by_pos(0.39, 3.16)
    if achievements_shape:
        desc = project.get('description_text', '')

        if desc:
            # Split description into sentences/lines for better bullet display
            # Limit to 3-4 bullet points for readability
            desc_lines = []
            # Split by periods or newlines
            for sentence in desc.replace('\n', '. ').split('. '):
                sentence = sentence.strip()
                if sentence and len(sentence) > 5:
                    # Truncate long sentences
                    if len(sentence) > 60:
                        sentence = sentence[:57] + '...'
                    desc_lines.append(sentence)
                    if len(desc_lines) >= 3:  # Max 3 bullet points
                        break

            if not desc_lines:
                desc_lines = [truncate_text(desc, 180)]

            set_shape_text_with_structure(
                achievements_shape,
                desc_lines,
                title_text=None,
                font_size=FONT_SIZES['content'],
                use_bullets=True
            )
        else:
            set_shape_text_with_structure(
                achievements_shape,
                ["No description available"],
                title_text=None,
                font_size=FONT_SIZES['content'],
                use_bullets=True
            )

    # =========================================================================
    # TRENDS (Shape 674) - CLEAR IT (overlaps with 679)
    # This shape overlaps with shape 679 (NEXT STEPS)
    #   - Shape 674 Y: 4.323" - 5.314"
    #   - Shape 679 Y: 4.532" - 5.423"
    # Progress info is now shown in SCOPE section (shape 672)
    # We must CLEAR this shape to avoid visual conflicts
    # =========================================================================
    trends_shape = find_shape_by_pos(0.39, 4.32)
    if trends_shape:
        set_shape_text(trends_shape, '', font_size=FONT_SIZES['content'])

    # =========================================================================
    # NEXT STEPS - User preference: Upcoming decisions
    # NOTE: Layout title "NEXT STEPS" is at (0.47, 4.14)
    #       Content shape at (0.39, 4.53) - 0.39" below title
    #       Template uses bullet formatting (▪)
    # =========================================================================
    next_shape = find_shape_by_pos(0.39, 4.53)
    if next_shape:
        next_items = []

        # Filter decisions that are not yet taken
        pending_decisions = [d for d in decisions if d.get('status') not in ['taken', 'actions-done']] if decisions else []

        if pending_decisions:
            for d in pending_decisions[:3]:
                title = d.get('title', d.get('name', ''))
                if title:
                    next_items.append(truncate_text(title, 45))
        elif decisions:
            # Show all decisions if none are pending
            for d in decisions[:3]:
                title = d.get('title', d.get('name', ''))
                if title:
                    next_items.append(truncate_text(title, 45))
        else:
            # Fallback to milestones if no decisions
            upcoming = [m for m in milestones if m.get('status') != 'done'] if milestones else []
            if upcoming:
                for m in upcoming[:3]:
                    name = m.get('name', '')
                    if name:
                        next_items.append(truncate_text(name, 45))

        if not next_items:
            next_items = ["No pending decisions or milestones"]

        set_shape_text_with_structure(
            next_shape,
            next_items,
            title_text=None,  # No inline title - layout has visual title above
            font_size=FONT_SIZES['content'],
            use_bullets=True
        )

    # =========================================================================
    # BUDGET - User preference: BAC + Actual + EAC
    # NOTE: Template has "Build" as first paragraph (title), then bullet items
    #       Layout title "BUDGET" is at (5.19, 3.96)
    # =========================================================================
    budget_shape = find_shape_by_pos(5.13, 4.4)
    if budget_shape:
        bac = project.get('budget_capex_initial')
        actual = project.get('budget_capex_used')
        eac = project.get('budget_capex_landing')

        budget_items = []
        budget_items.append(f"BAC: {bac:,.0f} €" if bac is not None else "BAC: N/A")
        budget_items.append(f"Actual: {actual:,.0f} €" if actual is not None else "Actual: N/A")
        budget_items.append(f"EAC: {eac:,.0f} €" if eac is not None else "EAC: N/A")

        set_shape_text_with_structure(
            budget_shape,
            budget_items,
            title_text="Build",  # Template has "Build" as inline title
            font_size=FONT_SIZES['content'],
            use_bullets=True
        )

    # =========================================================================
    # MADE - User preference: Milestones completados
    # NOTE: Template has "Made :" as first paragraph (title), then bullet items
    #       Layout title "DECISIONS" is at (5.19, 0.69) - but shape shows "Made"
    # =========================================================================
    made_shape = find_shape_by_pos(5.14, 1.0)
    if made_shape:
        completed = [m for m in milestones if m.get('status') == 'done'] if milestones else []

        made_items = []
        if completed:
            for m in completed[:4]:
                name = m.get('name', '')
                if name:
                    made_items.append(truncate_text(name, 40))

        if not made_items:
            made_items = ["No completed milestones yet"]

        set_shape_text_with_structure(
            made_shape,
            made_items,
            title_text="Made :",  # Template has "Made :" as inline title
            font_size=FONT_SIZES['content'],
            use_bullets=True
        )

    # =========================================================================
    # RISKS - User preference: Solo Risk Level (no attention points)
    # NOTE: Layout title "RISKS & STATUS" is at (5.19, 2.13)
    #       Template uses bullet formatting
    # =========================================================================
    risks_shape = find_shape_by_pos(5.14, 2.48)
    if risks_shape:
        risk_level = resolved.get('risk', project.get('risk', 'Not set'))

        risks_items = [f"Risk Level: {truncate_text(risk_level, 35)}"]

        # Optionally add attention points if they exist
        if attention_points:
            for ap in attention_points[:2]:  # Max 2 attention points
                ap_title = ap.get('title', '')
                if ap_title:
                    risks_items.append(truncate_text(ap_title, 40))

        set_shape_text_with_structure(
            risks_shape,
            risks_items,
            title_text=None,
            font_size=FONT_SIZES['content'],
            use_bullets=True
        )


def generate_from_template(data_path, output_path):
    """Generate PPT using the Systra template.

    This version properly preserves all template styling:
    - Images, borders, colors, backgrounds
    - Shape formatting and positioning
    - Fonts and text styles

    Strategy: Use template slide 0 as source, duplicate it for each project.
    The first project uses the original template slide, subsequent projects
    get copies of that slide with all its visual elements.

    Slide structure per spec section 5.1:
    1. Summary slide - List of all projects with mood/status
    2-N. Project slides - Cloned from template slide 1, with data filled
    Last. Data Notes slide - List of unfilled fields
    """
    from copy import deepcopy

    global UNFILLED_FIELDS
    UNFILLED_FIELDS = []  # Reset for new generation

    # TEMPLATE VERIFICATION (per spec: template is source of truth)
    print("Verifying template compatibility...")
    is_valid, warnings, detected_positions, full_report = verify_template(TEMPLATE_PATH)

    if warnings:
        for w in warnings:
            print(w)

    if not is_valid:
        print("\n❌ Template verification FAILED!")
        print("   The template structure has changed significantly.")
        print("   Run: python3 scripts/generate_ppt.py --analyze")
        print("   Then update EXPECTED_SHAPE_POSITIONS in the script.")
        print("\n   Proceeding anyway, but output may have misplaced content.\n")
    else:
        print("✓ Template verified OK\n")

    data = load_data(data_path)

    # Load template
    prs = Presentation(TEMPLATE_PATH)

    # Template has 3 slides:
    # 0: Project Review main (this is our template for project slides)
    # 1: Budget table
    # 2: Planning

    reference_data = data.get('reference_data', {})
    projects = data.get('projects', [])

    if not projects:
        print("No projects found in data file")
        return None

    print(f"Template loaded: {len(prs.slides)} slides")
    print(f"Processing {len(projects)} projects...")

    # ==========================================================================
    # EXTRACT ALL TEMPLATE STYLES
    # Before modifying any shapes, capture ALL template visual styles
    # (titles, content, bullets, dates, status) to apply uniformly
    # ==========================================================================
    global TEMPLATE_STYLES, TEMPLATE_BULLET_FORMAT
    TEMPLATE_STYLES['extracted'] = False  # Reset for new generation
    TEMPLATE_BULLET_FORMAT['extracted'] = False

    print("Extracting template styles...")
    template_slide = prs.slides[0]
    extract_all_template_styles(template_slide)

    # Delete the Budget and Planning template slides FIRST (index 1 and 2)
    # Do this before any other operations to avoid index confusion
    for slide_idx in [2, 1]:  # Delete in reverse order
        if slide_idx < len(prs.slides):
            rId = prs.slides._sldIdLst[slide_idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[slide_idx]

    # Now we have only slide 0 (Project Review template)
    # First, populate it with the first project data
    print(f"  Slide 1: {projects[0].get('project', {}).get('name', 'Unknown')}")
    populate_project_slide(prs.slides[0], projects[0], reference_data)
    track_project_unfilled(projects[0].get('project', {}))

    # For additional projects, duplicate slide 0
    for idx, proj_data in enumerate(projects[1:], start=1):
        project = proj_data.get('project', {})
        project_name = project.get('name', 'Unknown')
        print(f"  Slide {idx + 1}: {project_name}")

        # Duplicate the template slide (slide 0)
        new_slide = duplicate_slide(prs, 0)

        # Fill with project data
        populate_project_slide(new_slide, proj_data, reference_data)
        track_project_unfilled(project)

    # Add Summary slide at the end, then move to beginning
    print("Creating Summary slide...")
    create_summary_slide_on_template(prs, projects, reference_data)

    # Add Data Notes slide at the very end
    print("Creating Data Notes slide...")
    create_data_notes_slide(prs, UNFILLED_FIELDS)

    # Reorder: move Summary from second-to-last to position 0
    move_slide_to_position(prs, len(prs.slides) - 2, 0)

    # Save
    prs.save(output_path)
    print(f"\n✓ PPT saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  - Summary: 1")
    print(f"  - Project slides: {len(projects)}")
    print(f"  - Data Notes: 1")
    print(f"  Unfilled fields tracked: {len(UNFILLED_FIELDS)}")

    return output_path


def track_project_unfilled(project):
    """Track unfilled fields for a project."""
    project_name = project.get('name', 'Unknown')
    if not project.get('description_text'):
        track_unfilled_field(project_name, "Description", "Non renseigné")
    if not project.get('end_date'):
        track_unfilled_field(project_name, "End date", "Non renseigné")
    if project.get('budget_capex_initial') is None and project.get('budget_capex_used') is None:
        track_unfilled_field(project_name, "Budget", "Non renseigné dans AirSaas")


def move_slide_to_position(prs, from_index, to_index):
    """Move a slide from one position to another."""
    slides = prs.slides._sldIdLst
    slide_to_move = slides[from_index]

    # Remove from current position
    del slides[from_index]

    # Insert at new position
    slides.insert(to_index, slide_to_move)


def delete_placeholder_shapes(slide):
    """Delete all placeholder shapes from a slide.

    This completely removes the shapes rather than just clearing their text,
    which prevents "Double click to edit" from appearing.
    """
    shapes_to_delete = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            shapes_to_delete.append(shape)

    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)


def create_summary_slide_on_template(prs, projects, reference_data):
    """Create summary slide using a clean layout.

    Uses layout index 3 (Texte) which has minimal pre-existing shapes.
    """
    # Use "Texte" layout (index 3) - cleaner than Project_Card
    slide_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # DELETE all placeholder shapes completely (not just clear text)
    # This prevents "Double click to edit" from appearing
    delete_placeholder_shapes(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Portfolio Flash Report"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Subtitle with date
    subtitle_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.5), Inches(9), Inches(0.25))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Projets Vitaux CODIR - {datetime.now().strftime('%d/%m/%Y')}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Projects count
    count_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.8), Inches(9), Inches(0.25))
    tf = count_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{len(projects)} projets"
    p.font.size = Pt(10)
    p.font.italic = True

    # Table header - adjusted column positions for better alignment
    # Widened Status column and pushed Mood/Owner further right
    header_y = 1.1
    col_widths = [0.7, 2.8, 1.8, 1.5, 1.3]
    col_x = [0.3, 1.0, 3.8, 5.8, 7.5]

    headers = ["ID", "Projet", "Status", "Mood", "Owner"]
    for i, header in enumerate(headers):
        box = slide.shapes.add_textbox(Inches(col_x[i]), Inches(header_y), Inches(col_widths[i]), Inches(0.25))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Projects list
    row_y = header_y + 0.28
    row_height = 0.28

    for proj_data in projects[:14]:  # Can fit more rows now
        project = proj_data.get('project', {})
        resolved = proj_data.get('resolved', {})

        # Project ID
        short_id = project.get('short_id', '')
        box = slide.shapes.add_textbox(Inches(col_x[0]), Inches(row_y), Inches(col_widths[0]), Inches(row_height))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = short_id
        p.font.size = Pt(7)
        p.font.bold = True

        # Project Name
        name = project.get('name', 'Unknown')[:35]
        if len(project.get('name', '')) > 35:
            name += '...'
        box = slide.shapes.add_textbox(Inches(col_x[1]), Inches(row_y), Inches(col_widths[1]), Inches(row_height))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(7)

        # Status
        status = resolved.get('status', project.get('status', ''))
        box = slide.shapes.add_textbox(Inches(col_x[2]), Inches(row_y), Inches(col_widths[2]), Inches(row_height))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = status.replace('_', ' ').title() if status else '-'
        p.font.size = Pt(7)
        if status in STATUS_COLORS:
            p.font.color.rgb = STATUS_COLORS[status]

        # Mood with color
        mood = resolved.get('mood', project.get('mood', ''))
        mood_label = MOOD_LABELS.get(mood, mood) if mood else '-'
        box = slide.shapes.add_textbox(Inches(col_x[3]), Inches(row_y), Inches(col_widths[3]), Inches(row_height))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = mood_label
        p.font.size = Pt(7)
        if mood in MOOD_COLORS:
            p.font.color.rgb = MOOD_COLORS[mood]
            p.font.bold = True

        # Owner
        owner = project.get('owner', {})
        owner_name = owner.get('name', '') if isinstance(owner, dict) else ''
        if owner_name:
            parts = owner_name.split()
            owner_short = parts[0] if parts else ''
        else:
            owner_short = '-'
        box = slide.shapes.add_textbox(Inches(col_x[4]), Inches(row_y), Inches(col_widths[4]), Inches(row_height))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = owner_short
        p.font.size = Pt(7)

        row_y += row_height

    if len(projects) > 14:
        box = slide.shapes.add_textbox(Inches(0.3), Inches(row_y + 0.05), Inches(9), Inches(0.2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"... et {len(projects) - 14} autres projets"
        p.font.size = Pt(7)
        p.font.italic = True

    return slide


def update_planning_slide(slide, proj_data):
    """Update planning slide with milestone info."""
    milestones = proj_data.get('milestones', [])

    # Find the planning text shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            if 'Insert planning' in shape.text:
                if milestones:
                    ms_text = "Milestones:\n" + "\n".join([
                        f"{'✓' if m.get('status') == 'done' else '○'} {m.get('name', '')} - {m.get('date', '')}"
                        for m in milestones[:8]
                    ])
                else:
                    ms_text = "No milestones defined for this project"
                set_shape_text(shape, ms_text)
                break


def track_unfilled_field(project_name, field_name, reason):
    """Track unfilled fields for Data Notes slide."""
    global UNFILLED_FIELDS
    UNFILLED_FIELDS.append({
        'project': project_name,
        'field': field_name,
        'reason': reason
    })


def create_data_notes_slide(prs, unfilled_fields):
    """Create Data Notes slide listing unfilled fields.

    According to spec section 5.1:
    Last slide: Data Notes - List of unfilled fields
    """
    # Use "Texte" layout (index 3) - cleaner
    slide_layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # DELETE all placeholder shapes completely (not just clear text)
    # This prevents "Double click to edit" from appearing
    delete_placeholder_shapes(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Data Notes"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.5), Inches(9), Inches(0.25))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Champs non remplis / Fields not populated"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Known API limitations (always show these)
    api_limitations = [
        ("Mood comment", "L'API retourne uniquement le code du mood, pas le commentaire"),
        ("Deployment area", "Champ non disponible dans l'API AirSaas"),
        ("End users (actual/target)", "Champ non disponible dans l'API AirSaas"),
    ]

    # Section: API Limitations
    section_y = 0.85
    section_box = slide.shapes.add_textbox(Inches(0.4), Inches(section_y), Inches(9), Inches(0.25))
    tf = section_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Limitations API connues:"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x99, 0x33, 0x00)

    row_y = section_y + 0.28
    for field, reason in api_limitations:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(row_y), Inches(9), Inches(0.2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {field}: {reason}"
        p.font.size = Pt(7)
        row_y += 0.2

    # Section: Per-project unfilled fields
    if unfilled_fields:
        row_y += 0.15
        section_box = slide.shapes.add_textbox(Inches(0.4), Inches(row_y), Inches(9), Inches(0.25))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Champs manquants par projet:"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x99, 0x33, 0x00)

        row_y += 0.28

        # Group by project
        by_project = {}
        for item in unfilled_fields:
            proj = item['project']
            if proj not in by_project:
                by_project[proj] = []
            by_project[proj].append(f"{item['field']}: {item['reason']}")

        for project_name, fields in list(by_project.items())[:6]:  # Limit
            box = slide.shapes.add_textbox(Inches(0.5), Inches(row_y), Inches(9), Inches(0.18))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{project_name}:"
            p.font.size = Pt(7)
            p.font.bold = True
            row_y += 0.18

            for field_info in fields[:3]:
                box = slide.shapes.add_textbox(Inches(0.7), Inches(row_y), Inches(8.5), Inches(0.16))
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.text = f"- {field_info}"
                p.font.size = Pt(6)
                row_y += 0.16

            row_y += 0.08

    # Footer note
    footer_box = slide.shapes.add_textbox(Inches(0.4), Inches(5.2), Inches(9), Inches(0.2))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')} | Source: AirSaas API"
    p.font.size = Pt(6)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    return slide


def generate_ppt_simple(data_path, output_path):
    """Generate simple PPT without template (fallback).

    Slide structure per spec section 5.1:
    1. Summary slide - List of all projects with mood/status
    2-N. Project slides - Simple layout for each project
    Last. Data Notes slide - List of unfilled fields
    """
    global UNFILLED_FIELDS
    UNFILLED_FIELDS = []  # Reset for new generation

    data = load_data(data_path)

    # Create presentation (16:9)
    prs = Presentation()
    prs.slide_width = Inches(10)  # Match template
    prs.slide_height = Inches(5.625)

    projects = data.get('projects', [])
    reference_data = data.get('reference_data', {})

    # STEP 1: Create Summary slide first
    print("Creating Summary slide...")
    create_summary_slide(prs, projects, reference_data)

    # STEP 2: Add slides for each project
    print(f"Creating slides for {len(projects)} projects...")

    for proj_data in projects:
        project = proj_data.get('project', {})
        resolved = proj_data.get('resolved', {})
        project_name = project.get('name', 'Unknown')

        # Add a simple slide
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.17), Inches(8), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Project Review: {project_name}"
        p.font.size = Pt(20)
        p.font.bold = True

        # Content
        content_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.8), Inches(9), Inches(4))
        tf = content_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"Status: {resolved.get('status', project.get('status', ''))}"
        p.font.size = Pt(12)

        p = tf.add_paragraph()
        p.text = f"Mood: {resolved.get('mood', project.get('mood', ''))}"
        p.font.size = Pt(12)

        owner = project.get('owner', {})
        owner_name = owner.get('name', '') if isinstance(owner, dict) else ''
        p = tf.add_paragraph()
        p.text = f"Owner: {owner_name}"
        p.font.size = Pt(12)

        p = tf.add_paragraph()
        p.text = f"Timeline: {project.get('start_date', '')} → {project.get('end_date', '')}"
        p.font.size = Pt(12)

        desc = project.get('description_text', '')
        if desc:
            p = tf.add_paragraph()
            p.text = ""
            p = tf.add_paragraph()
            p.text = desc[:400] + ('...' if len(desc) > 400 else '')
            p.font.size = Pt(10)

        # Track unfilled fields
        if not desc:
            track_unfilled_field(project_name, "Description", "Non renseigné")
        if not project.get('end_date'):
            track_unfilled_field(project_name, "End date", "Non renseigné")
        if project.get('budget_capex_initial') is None and project.get('budget_capex_used') is None:
            track_unfilled_field(project_name, "Budget", "Non renseigné dans AirSaas")

    # STEP 3: Create Data Notes slide at the end
    print("Creating Data Notes slide...")
    create_data_notes_slide(prs, UNFILLED_FIELDS)

    prs.save(output_path)
    print(f"✓ PPT saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  - Summary: 1")
    print(f"  - Project slides: {len(projects)}")
    print(f"  - Data Notes: 1")
    print(f"  Unfilled fields tracked: {len(UNFILLED_FIELDS)}")

    return output_path


def generate_ppt(data_path, output_path):
    """Generate PPT - uses template if available, otherwise creates simple version."""
    if os.path.exists(TEMPLATE_PATH):
        print(f"Using template: {TEMPLATE_PATH}")
        return generate_from_template(data_path, output_path)
    else:
        print("Template not found, generating simple PPT")
        return generate_ppt_simple(data_path, output_path)


if __name__ == "__main__":
    import sys

    # Check for --analyze flag
    if len(sys.argv) > 1 and sys.argv[1] == '--analyze':
        verbose = '--verbose' in sys.argv or '-v' in sys.argv
        if os.path.exists(TEMPLATE_PATH):
            analyze_template(TEMPLATE_PATH, verbose=verbose)
        else:
            print(f"Error: Template not found at {TEMPLATE_PATH}")
        sys.exit(0)

    # Check for --export-mapping flag (exports current shape positions to JSON)
    if len(sys.argv) > 1 and sys.argv[1] == '--export-mapping':
        if os.path.exists(TEMPLATE_PATH):
            export_template_mapping(TEMPLATE_PATH)
        else:
            print(f"Error: Template not found at {TEMPLATE_PATH}")
        sys.exit(0)

    # Check for --verify flag
    if len(sys.argv) > 1 and sys.argv[1] == '--verify':
        if os.path.exists(TEMPLATE_PATH):
            print("Running template verification...\n")
            is_valid, warnings, positions, full_report = verify_template(TEMPLATE_PATH)
            for w in warnings:
                print(w)
            if is_valid:
                print("\n✓ Template is compatible")
            else:
                print("\n❌ Template has compatibility issues")

            # Show quick action items
            if full_report.get('drifted') or full_report.get('missing'):
                print("\n📋 ACTION ITEMS:")
                if full_report.get('drifted'):
                    print("   1. Update drifted shape positions in EXPECTED_SHAPE_POSITIONS")
                if full_report.get('missing'):
                    print("   2. Check if missing shapes exist in template or need to be removed from mapping")
                print("   3. Run --analyze for full shape details")

            sys.exit(0 if is_valid else 1)
        else:
            print(f"Error: Template not found at {TEMPLATE_PATH}")
            sys.exit(1)

    # Find latest data file
    data_dir = os.path.join(BASE_DIR, "data")
    data_files = sorted([f for f in os.listdir(data_dir) if f.endswith('_projects.json')], reverse=True)

    if not data_files:
        print("Error: No data files found. Run /fetch first.")
        sys.exit(1)

    data_path = os.path.join(data_dir, data_files[0])
    print(f"Using data: {data_files[0]}")

    # Output path
    today = datetime.now().strftime("%Y-%m-%d")
    output_path = os.path.join(BASE_DIR, "outputs", f"{today}_portfolio_skill.pptx")

    generate_ppt(data_path, output_path)
